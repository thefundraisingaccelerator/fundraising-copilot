import streamlit as st
from anthropic import Anthropic
import json
from pypdf import PdfReader
from pptx import Presentation
import tempfile
import os

# Page config
st.set_page_config(
    page_title="Fundraising Co-Pilot",
    page_icon="🚀",
    layout="centered",
    initial_sidebar_state="collapsed"
)

# Professional styling inspired by Claude's aesthetic
st.markdown("""
<style>
    /* Import clean fonts */
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600&family=Source+Serif+4:wght@400;500;600&display=swap');
    
    /* Hide Streamlit branding */
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    header {visibility: hidden;}
    
    /* Main container styling */
    .stApp {
        background-color: #FDFCFB;
    }
    
    .main .block-container {
        max-width: 720px;
        padding-top: 0.5rem;
        padding-bottom: 1rem;
    }
    
    /* Typography */
    h1, h2, h3, h4, h5, h6 {
        font-family: 'Source Serif 4', Georgia, serif !important;
        color: #1a1a1a;
    }
    
    p, span, div, input, textarea, button, label {
        font-family: 'Inter', -apple-system, BlinkMacSystemFont, sans-serif !important;
    }
    
    /* Header styling */
    .main-header {
        text-align: center;
        padding: 0.75rem 0 0.75rem 0;
        border-bottom: 1px solid #E8E4E0;
        margin-bottom: 1rem;
    }
    
    .main-header h1 {
        font-size: 1.75rem;
        font-weight: 600;
        color: #1a1a1a;
        margin-bottom: 0.4rem;
        letter-spacing: -0.02em;
    }
    
    .main-header .subtitle {
        font-size: 0.95rem;
        color: #666;
        font-weight: 400;
        line-height: 1.5;
        margin-bottom: 0.5rem;
    }
    
    .main-header .header-disclaimer {
        font-size: 0.75rem;
        color: #999;
        margin-top: 0.5rem;
    }
    
    /* Upload hint */
    .upload-hint {
        font-size: 0.85rem;
        color: #666;
        margin-bottom: 0.5rem;
        margin-top: 0.5rem;
    }
    
    /* Upload section */
    .upload-section {
        background: #FFFFFF;
        border: 1px solid #E8E4E0;
        border-radius: 12px;
        padding: 0.75rem;
        margin-bottom: 0.75rem;
    }
    
    .upload-label {
        font-size: 0.875rem;
        font-weight: 500;
        color: #1a1a1a;
        margin-bottom: 0.75rem;
        display: flex;
        align-items: center;
        gap: 0.5rem;
    }
    
    /* File uploader customization */
    .stFileUploader {
        background: transparent !important;
    }
    
    .stFileUploader > div {
        background: #FAFAFA !important;
        border: 1px dashed #D4D0CC !important;
        border-radius: 8px !important;
        padding: 1rem !important;
    }
    
    .stFileUploader > div:hover {
        border-color: #B8977E !important;
        background: #FDF9F6 !important;
    }
    
    /* Action buttons */
    .action-buttons {
        display: grid;
        grid-template-columns: 1fr 1fr;
        gap: 0.5rem;
        margin: 0.75rem 0;
    }
    
    .stButton > button {
        font-family: 'Inter', sans-serif !important;
        font-size: 0.9rem !important;
        font-weight: 500 !important;
        padding: 0.75rem 1rem !important;
        border-radius: 8px !important;
        border: 1px solid #E8E4E0 !important;
        background: #FFFFFF !important;
        color: #1a1a1a !important;
        transition: all 0.15s ease !important;
        width: 100% !important;
    }
    
    .stButton > button:hover {
        background: #FDF9F6 !important;
        border-color: #B8977E !important;
        color: #1a1a1a !important;
    }
    
    .stButton > button:active {
        background: #F5EDE6 !important;
    }
    
    /* Secondary button (clear chat) */
    .stButton > button[kind="secondary"] {
        background: transparent !important;
        border: 1px solid #E8E4E0 !important;
        color: #666 !important;
        font-size: 0.85rem !important;
        padding: 0.5rem 1rem !important;
    }
    
    .stButton > button[kind="secondary"]:hover {
        background: #F5F5F5 !important;
        color: #1a1a1a !important;
    }
    
    /* User message speech bubble */
    .user-message {
        background: #F0F0F0;
        border-radius: 18px;
        padding: 0.75rem 1rem;
        margin: 1rem 0;
        max-width: 85%;
        margin-left: auto;
        font-size: 0.95rem;
        line-height: 1.5;
        color: #1a1a1a;
    }
    
    /* Assistant response container */
    .assistant-container {
        margin: 1rem 0;
    }
    
    /* Make avatar image smaller and aligned left */
    .assistant-container img {
        border-radius: 50%;
        margin-bottom: 0.5rem;
    }
    
    /* Chat messages - hide default styling */
    .stChatMessage {
        background: transparent !important;
        border: none !important;
        padding: 1rem 0 !important;
    }
    
    .stChatMessage [data-testid="StyledLinkIconContainer"] {
        display: none !important;
    }
    
    /* User message */
    [data-testid="stChatMessageContent"]:has(> div > p) {
        font-size: 0.95rem !important;
        line-height: 1.6 !important;
    }
    
    /* Chat input */
    .stChatInput {
        border-top: 1px solid #E8E4E0;
        padding-top: 0.75rem;
        margin-top: 0.75rem;
    }
    
    .stChatInput > div {
        background: #FFFFFF !important;
        border: 1px solid #E8E4E0 !important;
        border-radius: 12px !important;
        padding: 0.25rem !important;
    }
    
    .stChatInput textarea {
        font-family: 'Inter', sans-serif !important;
        font-size: 0.95rem !important;
        color: #1a1a1a !important;
    }
    
    .stChatInput textarea::placeholder {
        color: #999 !important;
    }
    
    /* Success/info/warning messages */
    .stSuccess, .stInfo, .stWarning, .stError {
        font-size: 0.875rem !important;
        border-radius: 8px !important;
    }
    
    .stSuccess {
        background-color: #F0F9F4 !important;
        border: 1px solid #B8DBCA !important;
        color: #1a5d36 !important;
    }
    
    .stInfo {
        background-color: #F5F5F5 !important;
        border: 1px solid #E0E0E0 !important;
        color: #555 !important;
    }
    
    .stWarning {
        background-color: #FFF9F0 !important;
        border: 1px solid #F0D9B5 !important;
        color: #8B6914 !important;
    }
    
    /* Spinner */
    .stSpinner > div {
        border-color: #B8977E !important;
    }
    
    /* Footer */
    .footer {
        text-align: center;
        padding: 1rem 0 0.5rem 0;
        margin-top: 1rem;
        border-top: 1px solid #E8E4E0;
    }
    
    .footer p {
        font-size: 0.85rem;
        color: #888;
        margin: 0.2rem 0;
        line-height: 1.5;
    }
    
    .footer a {
        color: #B8977E;
        text-decoration: none;
        font-weight: 500;
    }
    
    .footer a:hover {
        color: #8B6B4A;
        text-decoration: underline;
    }
    
    .footer .tagline {
        font-family: 'Source Serif 4', Georgia, serif;
        font-style: italic;
        color: #666;
        font-size: 0.9rem;
        margin-top: 0.25rem;
    }
    
    /* Section divider */
    .section-divider {
        border: none;
        border-top: 1px solid #E8E4E0;
        margin: 1rem 0;
    }
    
    /* Prompt label */
    .prompt-label {
        font-size: 0.875rem;
        font-weight: 500;
        color: #555;
        margin-bottom: 0.5rem;
    }
    
    /* Hide default streamlit elements */
    .stDeployButton {display: none;}
    
    /* Hide empty containers */
    .stChatInput:empty, 
    .element-container:empty,
    .stMarkdown:empty {
        display: none !important;
    }
    
    /* Fix chat input container */
    [data-testid="stChatInput"] {
        background: transparent !important;
        border: none !important;
        box-shadow: none !important;
    }
    
    [data-testid="stChatInput"] > div {
        background: #FFFFFF !important;
        border: 1px solid #E8E4E0 !important;
        border-radius: 12px !important;
    }
    
    /* Sidebar styling */
    [data-testid="stSidebar"] {
        background: #FAFAFA !important;
        border-right: 1px solid #E8E4E0 !important;
    }
    
    [data-testid="stSidebar"] .stMarkdown {
        font-size: 0.875rem !important;
    }
    
    /* Markdown in responses */
    .stMarkdown {
        font-size: 0.95rem;
        line-height: 1.7;
        color: #1a1a1a;
    }
    
    .stMarkdown h3 {
        font-size: 1.1rem;
        margin-top: 1.5rem;
        margin-bottom: 0.75rem;
    }
    
    .stMarkdown ul, .stMarkdown ol {
        margin: 0.75rem 0;
        padding-left: 1.5rem;
    }
    
    .stMarkdown li {
        margin: 0.4rem 0;
    }
    
    .stMarkdown strong {
        font-weight: 600;
        color: #1a1a1a;
    }
    
    .stMarkdown code {
        background: #F5F5F5;
        padding: 0.15rem 0.4rem;
        border-radius: 4px;
        font-size: 0.85em;
    }
</style>
""", unsafe_allow_html=True)

# Load investor database
@st.cache_data
def load_investors():
    with open("investors.json", "r") as f:
        return json.load(f)

INVESTORS = load_investors()


def find_matching_investors(stage=None, sector_keywords=None, geography=None, investor_type=None, max_results=20):
    """Filter investors based on criteria"""
    matches = []
    
    for inv in INVESTORS:
        score = 0
        
        if stage and inv.get('stage'):
            inv_stages = inv['stage'].lower()
            if 'pre-seed' in stage.lower() or 'prototype' in stage.lower() or 'idea' in stage.lower():
                if 'prototype' in inv_stages or 'idea' in inv_stages or 'early revenue' in inv_stages:
                    score += 3
            elif 'seed' in stage.lower() or 'early revenue' in stage.lower():
                if 'early revenue' in inv_stages or 'prototype' in inv_stages:
                    score += 3
            elif 'series a' in stage.lower() or 'scaling' in stage.lower():
                if 'scaling' in inv_stages or 'growth' in inv_stages:
                    score += 3
        
        if sector_keywords and inv.get('thesis'):
            thesis_lower = inv['thesis'].lower()
            for keyword in sector_keywords:
                if keyword.lower() in thesis_lower:
                    score += 2
        
        if geography and inv.get('countries'):
            countries_lower = inv['countries'].lower()
            if geography.lower() in countries_lower or 'uk' in countries_lower:
                score += 2
        
        if investor_type:
            if investor_type.lower() in inv.get('type', '').lower():
                score += 1
        
        if score > 0 and (inv.get('thesis') or inv.get('stage')):
            matches.append((score, inv))
    
    matches.sort(key=lambda x: x[0], reverse=True)
    return [m[1] for m in matches[:max_results]]


def format_investor_for_context(investors):
    """Format investor list for inclusion in AI context"""
    if not investors:
        return "No matching investors found in the database."
    
    formatted = []
    for inv in investors:
        parts = [f"**{inv['name']}** ({inv['type']})"]
        if inv.get('stage'):
            parts.append(f"  - Stage: {inv['stage']}")
        if inv.get('thesis'):
            thesis = inv['thesis'][:300] + "..." if len(inv['thesis']) > 300 else inv['thesis']
            parts.append(f"  - Thesis: {thesis}")
        if inv.get('cheque_min') or inv.get('cheque_max'):
            cheque = f"{inv.get('cheque_min', '?')} - {inv.get('cheque_max', '?')}"
            parts.append(f"  - Cheque size: {cheque}")
        if inv.get('countries'):
            countries = inv['countries'][:100] + "..." if len(inv['countries']) > 100 else inv['countries']
            parts.append(f"  - Geography: {countries}")
        if inv.get('website'):
            parts.append(f"  - Website: {inv['website']}")
        formatted.append("\n".join(parts))
    
    return "\n\n".join(formatted)


def extract_text_from_pdf_basic(file):
    """Extract text from PDF using basic pypdf method"""
    reader = PdfReader(file)
    text = ""
    for page_num, page in enumerate(reader.pages, 1):
        page_text = page.extract_text() or ""
        if page_text.strip():
            text += f"\n--- Page {page_num} ---\n{page_text}"
    return text


def extract_text_from_pdf_ocr(file):
    """Extract text from PDF using OCR for image-heavy documents"""
    try:
        import pdf2image
        import pytesseract
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
            tmp_file.write(file.getvalue())
            tmp_path = tmp_file.name
        
        try:
            images = pdf2image.convert_from_path(tmp_path, dpi=150)
            
            text = ""
            for i, image in enumerate(images, 1):
                page_text = pytesseract.image_to_string(image)
                if page_text.strip():
                    text += f"\n--- Page {i} (OCR) ---\n{page_text}"
            
            return text
        finally:
            os.unlink(tmp_path)
            
    except ImportError as e:
        return None
    except Exception as e:
        st.warning(f"OCR processing error: {str(e)}")
        return None


def extract_text_from_pdf(file):
    """Extract text from PDF, trying basic extraction first then OCR if needed"""
    file.seek(0)
    basic_text = extract_text_from_pdf_basic(file)
    
    if basic_text and len(basic_text.strip()) > 500:
        return basic_text, "text"
    
    file.seek(0)
    ocr_text = extract_text_from_pdf_ocr(file)
    
    if ocr_text and len(ocr_text.strip()) > len(basic_text.strip() if basic_text else ""):
        return ocr_text, "OCR"
    
    return basic_text, "text"


def extract_text_from_pptx(file):
    """Extract text from PowerPoint file"""
    prs = Presentation(file)
    text = ""
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text += shape.text + "\n"
        if slide_text.strip():
            text += f"\n--- Slide {slide_num} ---\n{slide_text}"
    return text


def extract_deck_content(uploaded_file):
    """Extract text content from uploaded deck file"""
    if uploaded_file is None:
        return None, None
    
    try:
        if uploaded_file.type == "application/pdf":
            text, method = extract_text_from_pdf(uploaded_file)
            return text, method
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return extract_text_from_pptx(uploaded_file), "PPTX"
    except Exception as e:
        st.error(f"Error reading file: {str(e)}")
        return None, None
    
    return None, None


# System prompt
SYSTEM_PROMPT = """You are Fundraising Co-Pilot, an on-demand decision support assistant for early-stage founders who are actively fundraising or about to start.

Your role is to help founders make better fundraising decisions in real time, using an investor's perspective — so small mistakes don't compound.

## What You Help With
- Pressure-test pitch decks from an investor point of view
- Improve investor outreach emails before sending
- Sanity-check which investors are a realistic fit
- Clarify fundraising readiness and next priorities
- Understand likely objections investors will have
- Find relevant investors from the database

You explain WHY, not just what.

## Tone & Style
- Calm, direct, non-hypey
- Investor-realistic, not motivational
- Clear about trade-offs and uncertainty
- Assume the founder is smart but missing insider context
- Warm but honest - like a supportive mentor who tells hard truths

## Guardrails
You must never:
- Promise funding, responses, or introductions
- Claim certainty about investor decisions
- Act as legal, financial, or investment advice
- Encourage mass or untargeted investor outreach

If asked for guarantees: "There are no guarantees in fundraising — what I can do is help you reduce avoidable mistakes and improve clarity."

## How to Respond

**CRITICAL: When a pitch deck is provided, you MUST analyze THAT SPECIFIC DECK. Reference their actual content. Do not give generic advice.**

When a pitch deck is provided:
1. Start with a quick summary of what you understand the business to be
2. Identify the 2-3 biggest red flags an investor would notice
3. Point out what's unclear or missing
4. Give specific slide-by-slide observations where relevant
5. End with 2-3 priority fixes

When recommending investors:
1. Use the provided investor database matches
2. Explain why each investor might be a fit based on their thesis
3. Remind them to research each one and look for warm intro paths

## Key Heuristics

### Deck Red Flags
- **Vague Verbs**: "disrupting," "optimizing," "leveraging" without specifics
- **Mystery Product**: By slide 4, investor doesn't know what the product actually IS
- **Generic Titles**: "Our Solution" instead of "15% MoM Growth via Direct Sales"
- **Scale Mismatch**: Global problem → niche solution

### Pre-Seed Red Flags
- "We need money to build the MVP" (in 2026 with AI/no-code, this signals low resourcefulness)
- Core tech/sales outsourced to agency
- TAM = "1% of $100B market" (vs bottom-up: "5,000 law firms × £1k/mo")
- Messy cap table (advisors with 5% for nothing)

### What Investors Listen For
- **Earned Insight**: Non-obvious discovery from talking to 100+ customers
- **Speed of Iteration**: What happened since last meeting?
- **Unit Economics**: Know your CAC and margin
- **Why Now**: What changed recently that enables this?

### Too Early Signals
- Waitlist but zero pilots/LOIs
- Founder-problem mismatch (MedTech founder never worked in healthcare)
- Unclear what the money gets you to
- Only feedback from friends/family

## Default Framing
Use often: "From an investor's perspective…"

---
You are decision support, not a decision maker. Your goal is clarity, not confidence theatre.
"""

# Initialize Anthropic client
@st.cache_resource
def get_client():
    return Anthropic(api_key=st.secrets["ANTHROPIC_API_KEY"])

# Header with disclaimer
st.markdown("""
<div class="main-header">
    <h1>🚀 Fundraising Co-Pilot</h1>
    <p class="subtitle">AI-powered fundraising decision support for founders,<br>built with real knowledge and judgment from an investor with over 20+ years investing experience</p>
    <p class="header-disclaimer">Decision support, not advice. No guarantees in fundraising.</p>
</div>
""", unsafe_allow_html=True)

# Initialize session state
if "messages" not in st.session_state:
    st.session_state.messages = []
if "deck_content" not in st.session_state:
    st.session_state.deck_content = None
if "deck_filename" not in st.session_state:
    st.session_state.deck_filename = None

# Avatars for chat messages
ASSISTANT_AVATAR = "sutin_avatar.png"

# Display chat history
for message in st.session_state.messages:
    if message["role"] == "assistant":
        # Show avatar above the response, centered
        st.markdown('<div class="assistant-container">', unsafe_allow_html=True)
        st.image(ASSISTANT_AVATAR, width=36)
        st.markdown(message["content"])
        st.markdown('</div>', unsafe_allow_html=True)
    else:
        # User message as speech bubble (no avatar)
        st.markdown(f'<div class="user-message">{message["content"]}</div>', unsafe_allow_html=True)

# Starter prompts (only show if no messages)
if not st.session_state.messages:
    st.markdown('<p class="prompt-label">What can I help you with?</p>', unsafe_allow_html=True)
    
    col1, col2 = st.columns(2)
    with col1:
        if st.button("📊 Review my pitch deck", use_container_width=True):
            if st.session_state.deck_content:
                starter = "Review my pitch deck from an investor's perspective. Be specific about what's working and what needs to change."
            else:
                starter = "I'd like you to review my pitch deck. Let me upload it first."
            st.session_state.starter_prompt = starter
            st.rerun()
            
        if st.button("🎯 Am I ready to raise?", use_container_width=True):
            if st.session_state.deck_content:
                starter = "Based on my deck, am I ready to fundraise? What proof points am I missing?"
            else:
                starter = "Help me figure out if I'm ready to raise. What questions should I be able to answer before approaching investors?"
            st.session_state.starter_prompt = starter
            st.rerun()
            
    with col2:
        if st.button("🔍 Find investors for me", use_container_width=True):
            if st.session_state.deck_content:
                starter = "Based on my pitch deck, find investors who would be a good fit for my startup."
            else:
                starter = "Help me find investors. I'll describe my startup so you can suggest who might be a good fit."
            st.session_state.starter_prompt = starter
            st.rerun()
            
        if st.button("✉️ Review my outreach email", use_container_width=True):
            if st.session_state.deck_content:
                starter = "Help me write a cold email to an investor based on my deck. What should I include to get a response?"
            else:
                starter = "I want to write a cold email to an investor. What makes the difference between one that gets ignored vs one that gets a response?"
            st.session_state.starter_prompt = starter
            st.rerun()
    
    # Upload section - simple and clean, no expander
    st.markdown('<p class="upload-hint">Have a pitch deck? Upload it for specific feedback.</p>', unsafe_allow_html=True)
    uploaded_file = st.file_uploader(
        "Upload deck",
        type=["pdf", "pptx"],
        help="PDF or PowerPoint. We'll extract the text to give you specific feedback.",
        label_visibility="collapsed"
    )
    
    # Process uploaded file
    if uploaded_file is not None:
        if st.session_state.deck_filename != uploaded_file.name:
            with st.spinner("Processing your deck..."):
                deck_content, method = extract_deck_content(uploaded_file)
                
            if deck_content and len(deck_content.strip()) > 100:
                st.session_state.deck_content = deck_content
                st.session_state.deck_filename = uploaded_file.name
                st.success(f"✓ Ready: {uploaded_file.name}")
            else:
                st.error("Couldn't extract content. Try a different file.")
        else:
            st.success(f"✓ Using: {uploaded_file.name}")

else:
    # When in conversation, show smaller upload option if no deck loaded
    if not st.session_state.deck_content:
        with st.sidebar:
            st.markdown("**📎 Add your deck**")
            uploaded_file = st.file_uploader(
                "Upload",
                type=["pdf", "pptx"],
                label_visibility="collapsed",
                key="sidebar_upload"
            )
            if uploaded_file is not None:
                if st.session_state.deck_filename != uploaded_file.name:
                    with st.spinner("Processing..."):
                        deck_content, method = extract_deck_content(uploaded_file)
                    if deck_content and len(deck_content.strip()) > 100:
                        st.session_state.deck_content = deck_content
                        st.session_state.deck_filename = uploaded_file.name
                        st.success(f"✓ {uploaded_file.name}")
                        st.rerun()
    else:
        with st.sidebar:
            st.markdown(f"**📄 Deck loaded**")
            st.caption(st.session_state.deck_filename)
            if st.button("Remove", type="secondary"):
                st.session_state.deck_content = None
                st.session_state.deck_filename = None
                st.rerun()

# Handle starter prompts
if "starter_prompt" in st.session_state:
    prompt = st.session_state.starter_prompt
    del st.session_state.starter_prompt
    
    st.session_state.messages.append({"role": "user", "content": prompt})
    
    # Check if this is an investor search request
    investor_keywords = ['find investor', 'find investors', 'suggest investor', 'recommend investor', 'who should i pitch', 'match my startup']
    is_investor_search = any(kw in prompt.lower() for kw in investor_keywords)
    
    # Check if this is a deck review request
    deck_review_keywords = ['review my pitch deck', 'review my deck', 'analyze my deck', 'feedback on my deck']
    is_deck_review = any(kw in prompt.lower() for kw in deck_review_keywords)
    
    full_prompt = prompt
    additional_context = ""
    
    # SCENARIO 1: Investor search WITH deck - search database and recommend
    if is_investor_search and st.session_state.deck_content:
        additional_context += f"""

---
**PITCH DECK CONTENT** (from {st.session_state.deck_filename}):

{st.session_state.deck_content[:15000]}

---
Use this deck to understand the business and find matching investors.
"""
        # Extract stage and sectors from deck
        stage = None
        sector_keywords = []
        sectors = ['ai', 'fintech', 'healthtech', 'health', 'saas', 'b2b', 'b2c', 'consumer', 'enterprise', 
                   'climate', 'sustainability', 'edtech', 'proptech', 'foodtech', 'biotech', 'deeptech',
                   'marketplace', 'ecommerce', 'gaming', 'web3', 'blockchain', 'crypto', 'mental health', 
                   'wellness', 'fashion', 'retail', 'logistics', 'hr', 'legal', 'insurance', 'cybersecurity', 
                   'iot', 'robotics', 'energy', 'cleantech', 'agtech', 'space', 'mobility', 'impact',
                   'neurodiversity', 'diversity', 'inclusion', 'workplace', 'employee', 'future of work']
        
        search_text = st.session_state.deck_content.lower()
        
        # Detect stage from deck
        if any(s in search_text for s in ['pre-seed', 'preseed', 'idea stage', 'prototype']):
            stage = "pre-seed"
        elif any(s in search_text for s in ['seed', 'early revenue', 'mvp', 'pilot', 'first customer']):
            stage = "seed"
        elif any(s in search_text for s in ['series a', 'scaling', 'growth stage']):
            stage = "series a"
        else:
            stage = "seed"  # Default for most decks
        
        for sector in sectors:
            if sector in search_text:
                sector_keywords.append(sector)
        
        matches = find_matching_investors(
            stage=stage,
            sector_keywords=sector_keywords[:5] if sector_keywords else None,
            geography="UK",
            max_results=15
        )
        
        if matches:
            additional_context += f"""

---
**MATCHING INVESTORS FROM DATABASE**:

{format_investor_for_context(matches)}

Based on the deck and these investor matches, recommend 5-10 investors that fit best. Explain why each is a good fit based on their thesis and the startup's focus. Remind them to research each one and look for warm intro paths.
---
"""
    
    # SCENARIO 2: Investor search WITHOUT deck - ask for details
    elif is_investor_search and not st.session_state.deck_content:
        additional_context += """

The user wants help finding investors but hasn't uploaded a deck or described their startup yet. 
Ask them to briefly describe: 1) What their startup does, 2) What stage they're at, 3) What sector/industry they're in.
Once they provide this, you can search the investor database for matches.
"""
    
    # SCENARIO 3: Deck review WITH deck - analyze it
    elif is_deck_review and st.session_state.deck_content:
        additional_context += f"""

---
**PITCH DECK CONTENT** (from {st.session_state.deck_filename}):

{st.session_state.deck_content[:15000]}

---
Analyze THIS SPECIFIC DECK. Reference their actual slides and content. Do not give generic advice.
"""
    
    # SCENARIO 4: Deck review WITHOUT deck - prompt to upload
    elif is_deck_review and not st.session_state.deck_content:
        additional_context += """

The user wants a deck review but hasn't uploaded one yet. Let them know they can upload a PDF or PowerPoint deck using the file uploader, and you'll give specific feedback on it.
"""
    
    # SCENARIO 5: Other requests WITH deck - reference it where relevant  
    elif st.session_state.deck_content:
        additional_context += f"""

---
**PITCH DECK CONTENT** (from {st.session_state.deck_filename}):

{st.session_state.deck_content[:15000]}

---
Reference this deck content in your response where relevant.
"""
    
    # SCENARIO 6: Other requests WITHOUT deck - just respond normally
    # (no additional context needed)
    
    full_prompt = prompt + additional_context
    
    client = get_client()
    
    # Show avatar above response
    st.markdown('<div class="assistant-container">', unsafe_allow_html=True)
    st.image(ASSISTANT_AVATAR, width=36)
    with st.spinner("Analyzing..."):
        response = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=2500,
            system=SYSTEM_PROMPT,
            messages=[{"role": "user", "content": full_prompt}]
        )
        assistant_message = response.content[0].text
    st.markdown(assistant_message)
    st.markdown('</div>', unsafe_allow_html=True)
    
    st.session_state.messages.append({"role": "assistant", "content": assistant_message})
    st.rerun()

# Chat input
if prompt := st.chat_input("Ask a fundraising question..."):
    st.session_state.messages.append({"role": "user", "content": prompt})
    
    # Display user message as speech bubble
    st.markdown(f'<div class="user-message">{prompt}</div>', unsafe_allow_html=True)
    
    # Detect intent
    investor_keywords = ['find investor', 'find investors', 'suggest investor', 'recommend investor', 'who should i pitch', 'match', 'which vc', 'which angel']
    is_investor_search = any(kw in prompt.lower() for kw in investor_keywords)
    
    additional_context = ""
    
    # Add deck content if available
    if st.session_state.deck_content:
        additional_context += f"""

---
**PITCH DECK CONTENT** (from {st.session_state.deck_filename}):

{st.session_state.deck_content[:15000]}

---
Reference this deck content in your response where relevant.
"""
    
    # If investor search, try to find matches
    if is_investor_search:
        # Build search text from prompt + deck (if available) + conversation history
        search_text = prompt.lower()
        
        # Add deck content to search
        if st.session_state.deck_content:
            search_text += " " + st.session_state.deck_content.lower()
        
        # Also check recent conversation for context (user might have described their startup)
        for msg in st.session_state.messages[-6:]:  # Last few messages
            if msg["role"] == "user":
                search_text += " " + msg["content"].lower()
        
        # Detect stage
        stage = None
        if any(s in search_text for s in ['pre-seed', 'preseed', 'idea stage', 'prototype']):
            stage = "pre-seed"
        elif any(s in search_text for s in ['seed', 'early revenue', 'mvp', 'pilot', 'first customer']):
            stage = "seed"
        elif any(s in search_text for s in ['series a', 'scaling', 'growth stage']):
            stage = "series a"
        
        # Detect sectors
        sector_keywords = []
        sectors = ['ai', 'fintech', 'healthtech', 'health', 'saas', 'b2b', 'b2c', 'consumer', 'enterprise', 
                   'climate', 'sustainability', 'edtech', 'proptech', 'foodtech', 'biotech', 'deeptech',
                   'marketplace', 'ecommerce', 'gaming', 'web3', 'blockchain', 'crypto', 'mental health', 
                   'wellness', 'fashion', 'retail', 'logistics', 'hr', 'legal', 'insurance', 'cybersecurity', 
                   'iot', 'robotics', 'energy', 'cleantech', 'agtech', 'space', 'mobility', 'impact',
                   'neurodiversity', 'diversity', 'inclusion', 'workplace', 'employee', 'future of work']
        
        for sector in sectors:
            if sector in search_text:
                sector_keywords.append(sector)
        
        # Detect geography
        geography = None
        if any(g in prompt.lower() for g in ['uk', 'united kingdom', 'london', 'britain']):
            geography = "UK"
        elif any(g in prompt.lower() for g in ['us', 'usa', 'united states', 'america']):
            geography = "USA"
        elif any(g in prompt.lower() for g in ['europe', 'eu']):
            geography = "Europe"
        else:
            geography = "UK"  # Default
        
        # Only search if we have enough context
        if stage or sector_keywords:
            matches = find_matching_investors(
                stage=stage,
                sector_keywords=sector_keywords[:5] if sector_keywords else None,
                geography=geography,
                max_results=15
            )
            
            if matches:
                additional_context += f"""

---
**MATCHING INVESTORS FROM DATABASE**:

{format_investor_for_context(matches)}

Recommend 5-10 that fit best based on the startup described. Explain why each is a good fit. Remind them to research and find warm intros.
---
"""
        else:
            # Not enough info to search - prompt will ask for details
            additional_context += """

The user wants investor recommendations but you don't have enough context about their startup yet. 
Ask them to describe: 1) What their startup does, 2) What stage they're at (pre-seed, seed, Series A), 3) What sector/industry.
"""
    
    client = get_client()
    
    # Show avatar above response
    st.markdown('<div class="assistant-container">', unsafe_allow_html=True)
    st.image(ASSISTANT_AVATAR, width=36)
    with st.spinner(""):
        messages_for_api = [{"role": m["role"], "content": m["content"]} for m in st.session_state.messages[:-1]]
        messages_for_api.append({"role": "user", "content": prompt + additional_context})
        
        response = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=2500,
            system=SYSTEM_PROMPT,
            messages=messages_for_api
        )
        assistant_message = response.content[0].text
    st.markdown(assistant_message)
    st.markdown('</div>', unsafe_allow_html=True)
    
    st.session_state.messages.append({"role": "assistant", "content": assistant_message})

# Clear button (only show if there are messages)
if st.session_state.messages:
    st.markdown("<br>", unsafe_allow_html=True)
    if st.button("↻ Start over", type="secondary"):
        st.session_state.messages = []
        st.session_state.deck_content = None
        st.session_state.deck_filename = None
        st.rerun()

# Footer
st.markdown("""
<div class="footer">
    <p>Built by <a href="https://thefundraisingaccelerator.com" target="_blank">The Fundraising Accelerator</a></p>
    <p class="tagline">Your network should not determine your net worth.</p>
</div>
""", unsafe_allow_html=True)
